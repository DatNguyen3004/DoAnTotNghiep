import collections
from pptx import Presentation
import sys

def read_pptx(filepath, outpath):
    try:
        prs = Presentation(filepath)
        with open(outpath, 'w', encoding='utf-8') as f:
            f.write(f"File: {filepath}\n")
            f.write(f"Total Slides: {len(prs.slides)}\n")
            for i, slide in enumerate(prs.slides):
                f.write(f"--- Slide {i+1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        f.write(shape.text.strip() + "\n")
                f.write("\n")
    except Exception as e:
        print(f"Error reading {filepath}: {e}")

if __name__ == "__main__":
    read_pptx(sys.argv[1], sys.argv[2])
