from pptx import Presentation
import sys

def inspect_slide_3(filepath):
    prs = Presentation(filepath)
    slide = prs.slides[2] # Slide 3 is index 2
    print("--- Slide 3 Shapes ---")
    for i, shape in enumerate(slide.shapes):
        print(f"Shape {i}: Type: {shape.shape_type}")
        if hasattr(shape, "text") and shape.text.strip():
            print(f"  Text: {repr(shape.text)}")
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    print(f"    Paragraph: {repr(paragraph.text)}")

if __name__ == "__main__":
    inspect_slide_3("DATN/Đồ án Tốt nghiệp - NuLabel(1).pptx")
